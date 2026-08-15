from __future__ import annotations

import re
import subprocess
from pathlib import Path

from .core import ValidationError


def _bounds(spec: str | None, prefix: str, length: int) -> tuple[int, int]:
    if not spec:
        return 0, length
    match = re.fullmatch(rf"{prefix}(\d+)-(\d+)", spec)
    if not match:
        raise ValidationError(f"invalid slice {spec!r}; expected {prefix}<start>-<end>")
    start, end = map(int, match.groups())
    if start < 1 or end < start:
        raise ValidationError(f"invalid slice: {spec}")
    return start - 1, min(end, length)


def render_file(path: Path, slice_spec: str | None = None) -> str:
    suffix = path.suffix.lower()
    if suffix in {".txt", ".md", ".vtt", ".srt"}:
        text = path.read_text(encoding="utf-8", errors="replace")
        if slice_spec and slice_spec.startswith("h2:"):
            heading = slice_spec[3:]
            match = re.search(rf"(?mi)^##\s+{re.escape(heading)}\s*$", text)
            if not match:
                raise ValidationError(f"heading not found: {heading}")
            following = re.search(r"(?m)^#{1,2}\s+", text[match.end() :])
            return text[match.start() : match.end() + (following.start() if following else len(text))].strip() + "\n"
        return text
    if suffix == ".pdf":
        try:
            from pypdf import PdfReader

            pages = PdfReader(str(path)).pages
            start, end = _bounds(slice_spec, "p", len(pages))
            return "\n\n".join(f"## Page {i + 1}\n\n{pages[i].extract_text() or ''}" for i in range(start, end)) + "\n"
        except ImportError as exc:
            raise ValidationError("PDF rendering requires pypdf") from exc
    if suffix == ".pptx":
        try:
            from pptx import Presentation

            slides = Presentation(str(path)).slides
            start, end = _bounds(slice_spec, "s", len(slides))
            chunks = []
            for index in range(start, end):
                slide = slides[index]
                text = [shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text]
                notes = getattr(slide, "notes_slide", None)
                if notes:
                    note_text = [s.text for s in notes.notes_text_frame.paragraphs if getattr(s, "text", "")]
                    if note_text:
                        text.extend(["### Speaker notes", *note_text])
                chunks.append(f"## Slide {index + 1}\n\n" + "\n\n".join(text))
            return "\n\n".join(chunks) + "\n"
        except ImportError as exc:
            raise ValidationError("PowerPoint rendering requires python-pptx") from exc
    if suffix == ".docx":
        try:
            from docx import Document

            return "\n\n".join(p.text for p in Document(str(path)).paragraphs if p.text) + "\n"
        except ImportError as exc:
            raise ValidationError("Word rendering requires python-docx") from exc
    raise ValidationError(
        f"no converter for {suffix or path.name}; implement the converter interface in pruefung.converters"
    )


def render_url(url: str) -> tuple[bytes, str]:
    try:
        import httpx

        response = httpx.get(url, follow_redirects=True, timeout=30)
        response.raise_for_status()
    except Exception as exc:
        raise ValidationError(f"could not fetch {url}: {exc}") from exc
    raw = response.content
    text = re.sub(r"(?is)<(script|style).*?>.*?</\1>", "", response.text)
    text = re.sub(r"(?s)<[^>]+>", " ", text)
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n\s*\n+", "\n\n", text)
    return raw, text.strip() + "\n"


def render_git(locator: str, rev: str | None, paths: str | None) -> str:
    command = ["git", "-C", locator, "log", "-p", "--no-color"]
    if rev:
        command.append(rev)
    if paths:
        command.extend(["--", paths])
    result = subprocess.run(command, capture_output=True, text=True, check=False)
    if result.returncode:
        raise ValidationError(result.stderr.strip() or "git render failed")
    return "# Git history\n\n```diff\n" + result.stdout + "\n```\n"
